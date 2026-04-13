import os
import logging
import time
import requests
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from core.config import Config
from core.llm_tracker import track as _track

logger = logging.getLogger(__name__)


class MarketingAgent:
    THEMES = {
        "TECH": {"bg": (15, 23, 42), "title": (56, 189, 248), "text": (226, 232, 240)}, # Slate-900 / Sky-400 / Slate-200
        "ECO": {"bg": (20, 83, 45), "title": (74, 222, 128), "text": (220, 252, 231)},   # Green-900 / Green-400 / Green-100
        "CORPORATE": {"bg": (255, 255, 255), "title": (30, 58, 138), "text": (71, 85, 105)}, # White / Blue-900 / Slate-600
        "LUXURY": {"bg": (0, 0, 0), "title": (212, 175, 55), "text": (250, 250, 250)},   # Black / Gold / White
        "CREATIVE": {"bg": (255, 241, 242), "title": (190, 18, 60), "text": (88, 28, 135)} # Pink-50 / Rose-700 / Purple-900
    }

    def __init__(self):
        from core.clients import get_openai_client, get_runway_client
        self.openai = get_openai_client()
        self.runway_client = get_runway_client()
        self.user_phone: str | None = None
        self.base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        self.output_folder = os.path.join(self.base_dir, 'static', 'generated_docs')
        os.makedirs(self.output_folder, exist_ok=True)

    def generate_marketing_content(self, prompt_text, content_type="image", business_name="Mi Empresa", logo_path=None, user_phone=None):
        """
        Genera contenido creativo.
        content_type: 'image' (DALL-E 3) o 'slide' (PowerPoint).
        """
        self.user_phone = user_phone
        logger.info("MarketingAgent: creando %s para %s", content_type, business_name)

        if content_type == "image":
            return self._generate_dalle_image(prompt_text, business_name, logo_path)
        elif content_type == "slide":
            return self._generate_pptx_slide(prompt_text, business_name)
        elif content_type == "video":
            # 1. Visual Intelligence: Analizar Imagen para crear Prompt
            target_image = logo_path if logo_path else "dummy_product.jpg"
            video_prompt = self._analyze_product_context(target_image, business_name=business_name)

            # 2. Generar Video con Runway (Image + Text)
            # Pasamos la imagen original y el prompt mejorado
            video_url = self._generate_runway_video(video_prompt, input_image_path=target_image)
            return video_url
        else:
            return None

    def _generate_dalle_image(self, user_prompt, business_name, logo_path=None):
        try:
            # 1. Optimizar prompt para DALL-E usando GPT-4o
            optimized_prompt = self._refine_prompt_for_dalle(user_prompt, business_name)
            logger.info("Prompt DALL-E: %s...", optimized_prompt[:50])
            
            # 2. Generar Imagen
            response = self.openai.images.generate(
                model="dall-e-3",
                prompt=optimized_prompt,
                size="1024x1024",
                quality="standard",
                n=1,
            )
            image_url = response.data[0].url
            
            # 3. Descargar y guardar localmente
            filename = f"marketing_{int(datetime.now().timestamp())}.png"
            local_path = os.path.join(self.output_folder, filename)
            
            img_data = requests.get(image_url).content
            with open(local_path, 'wb') as handler:
                handler.write(img_data)
                
            # 4. Superponer Logo (Si existe)
            if logo_path:
                try:
                    from PIL import Image
                    # Convertir ruta web relativa a absoluta
                    # logo_path suele ser "/static/uploads/logos/..."
                    # self.base_dir está en "modules/proactive/...", subir 3 niveles para llegar a root
                    # Pero self.output_folder ya usa self.base_dir correctamente calculado en __init__
                    # self.base_dir en __init__ es: os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
                    # Esto apunta a TICKETIA_PRO root.
                    # logo_path empieza con /, quitarlo.
                    abs_logo_path = os.path.join(self.base_dir, logo_path.lstrip('/'))
                    
                    if os.path.exists(abs_logo_path):
                        base_img = Image.open(local_path).convert("RGBA")
                        logo_img = Image.open(abs_logo_path).convert("RGBA")
                        
                        # Resize logo (ej: 20% del ancho)
                        base_w, base_h = base_img.size
                        target_w = int(base_w * 0.20)
                        aspect_ratio = logo_img.width / logo_img.height
                        target_h = int(target_w / aspect_ratio)
                        logo_img = logo_img.resize((target_w, target_h), Image.Resampling.LANCZOS)
                        
                        # Posición: Esquina inferior derecha con padding
                        padding = 30
                        pos_x = base_w - target_w - padding
                        pos_y = base_h - target_h - padding
                        
                        # Pegar (usando el mismo logo como máscara para transparencia)
                        base_img.paste(logo_img, (pos_x, pos_y), logo_img)
                        
                        # Guardar (sobreescribir)
                        base_img.save(local_path, format="PNG")
                        logger.info("Logo superpuesto correctamente")
                    else:
                        logger.warning("Logo no encontrado: %s", abs_logo_path)
                except Exception as e_pil:
                    logger.warning("Error procesando logo PIL: %s", e_pil)
                
            return f"{Config.PUBLIC_URL}/static/generated_docs/{filename}"
            
        except Exception as e:
            logger.error("Error DALL-E: %s", e)
            return None

    def _analyze_product_context(self, image_path, business_name=""):
        """
        Two-stage Visual Intelligence pipeline:
        Stage 1 – GPT-4o Vision: understand WHAT is in the image (product type, colors, natural context).
        Stage 2 – GPT-4o Text: convert that analysis into a rich cinematic Runway Gen-3 prompt.
        """
        import base64
        import mimetypes

        if not os.path.exists(image_path):
            logger.warning("Imagen no encontrada para analisis: %s", image_path)
            return "Cinematic tracking shot of a person using the product outdoors, natural lighting, movement, photorealistic, 4K"

        mime_type, _ = mimetypes.guess_type(image_path)
        if not mime_type:
            mime_type = "image/png"

        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
            base64_image = f"data:{mime_type};base64,{encoded_string}"

        try:
            # ── Stage 1: Image Understanding ──────────────────────────────────
            logger.info("Stage 1 – Analizando imagen: %s", os.path.basename(image_path))

            _t0 = time.time()
            stage1_response = self.openai.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are a product analyst. Describe the image in structured detail:\n"
                            "1. PRODUCT_TYPE: What is the main subject/product?\n"
                            "2. CATEGORY: clothing / food / electronics / tool / beauty / decor / other\n"
                            "3. COLORS: Main colors and materials\n"
                            "4. STYLE: Style, mood, target audience\n"
                            "5. NATURAL_CONTEXT: Where would this product realistically be used, worn, or shown?\n"
                            "Be concise. Max 150 words."
                        )
                    },
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Analyze this product image."},
                            {
                                "type": "image_url",
                                "image_url": {"url": base64_image, "detail": "high"}
                            }
                        ]
                    }
                ],
                max_tokens=200
            )
            _track(self.user_phone, "gpt-4o", "video_analyze_image",
                   stage1_response, int((time.time() - _t0) * 1000))

            product_analysis = stage1_response.choices[0].message.content.strip()
            logger.info("Análisis del producto:\n%s", product_analysis)

            # ── Stage 2: Cinematic Prompt Generation ──────────────────────────
            logger.info("Stage 2 – Generando prompt cinematográfico")

            stage2_system = """You are a world-class Runway Gen-3 Alpha prompt engineer specializing in commercial video ads.

Your job: Given a product analysis, write a single cinematic Runway prompt that forces the AI to generate REAL MOVEMENT — never a static slideshow.

RUNWAY PROMPT FORMULA:
[CAMERA MOVEMENT] of [HUMAN/SUBJECT IN ACTION] [WEARING/USING/INTERACTING WITH the product] in [VIVID ENVIRONMENT]. [SPECIFIC MOTION DETAILS]. [LIGHTING]. [MOOD]. Cinematic, photorealistic, 4K.

CAMERA MOVEMENTS (pick the most fitting):
"Slow dolly-in" / "Tracking shot" / "Low-angle push-in" / "Aerial drone" / "Close-up slow-motion" / "Handheld follow"

MOTION RULES (CRITICAL):
- ALWAYS include a human or living subject actively interacting with the product
- ALWAYS describe movement: walking, pouring, cutting, spinning, flowing, jumping
- ALWAYS include environmental motion: wind, water, steam, particles, bokeh, reflections
- NEVER describe the product as static, isolated, or floating

EXAMPLES:
- Clothing: "Tracking shot of a fashion model wearing a white linen beach shirt, walking barefoot on wet sand at golden hour, fabric flowing in the ocean breeze, warm sunlight casting long shadows, shallow depth of field. Cinematic, photorealistic, 4K."
- Food: "Slow-motion close-up of a silver knife slicing through a layered chocolate cake, creamy filling revealed, crumbs scattering in the air, warm ambient kitchen lighting, steam rising. Cinematic, 4K."
- Electronics: "Low-angle push-in shot of a designer's hands using a sleek laptop on a rooftop terrace at dusk, screen glow on face, city lights bokeh in background, cool blue tones. Cinematic, 4K."
- Tool: "Handheld follow shot of a craftsman using a power drill to bore into a wooden wall, sawdust flying, focused expression, workshop environment with dramatic side lighting. Cinematic, 4K."

OUTPUT: Only the final prompt. No explanations. In English. Max 80 words."""

            brand_context = f" The brand is '{business_name}'." if business_name else ""

            _t0 = time.time()
            stage2_response = self.openai.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": stage2_system},
                    {
                        "role": "user",
                        "content": (
                            f"Product analysis:{brand_context}\n\n{product_analysis}\n\n"
                            "Generate the perfect Runway video prompt."
                        )
                    }
                ],
                max_tokens=350
            )
            _track(self.user_phone, "gpt-4o", "video_generate_prompt",
                   stage2_response, int((time.time() - _t0) * 1000))

            final_prompt = stage2_response.choices[0].message.content.strip()
            logger.info("Prompt final para Runway: %s", final_prompt)
            return final_prompt

        except Exception as e:
            logger.error("Error Visual Intelligence: %s", e)
            return "Cinematic tracking shot of a person using the product outdoors, natural lighting, wind, movement, photorealistic, 4K"
    def _generate_pptx_slide(self, topic, business_name):
        try:
            # 1. Estructurar contenido COMPLETO + TEMA
            presentation_data = self._get_presentation_structure(topic)
            theme_key = presentation_data.get('theme', 'CORPORATE').upper()
            theme = self.THEMES.get(theme_key, self.THEMES['CORPORATE'])
            
            # 2. Crear PPT
            prs = Presentation()
            
            # --- Slide 1: Título Principal ---
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            self._apply_theme(slide, theme) # Aplicar estilos
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = presentation_data.get('title', 'Propuesta')
            # Incluir nombre empresa en subtítulo
            subtitle.text = f"{presentation_data.get('subtitle', '')}\nPropuesto por: {business_name}"
            
            self._style_text(title, theme['title'], bold=True, size=Pt(44))
            self._style_text(subtitle, theme['text'], size=Pt(24))

            # --- Slides de Contenido ---
            for slide_data in presentation_data.get('slides', []):
                bullet_slide_layout = prs.slide_layouts[1] # Título + Bullets
                slide = prs.slides.add_slide(bullet_slide_layout)
                self._apply_theme(slide, theme)
                
                # Título de la Slide
                shapes = slide.shapes
                title_shape = shapes.title
                title_shape.text = slide_data.get('title', 'Detalle')
                self._style_text(title_shape, theme['title'], bold=True, size=Pt(36))
                
                # Cuerpo (Bullets)
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.word_wrap = True
                # Limpiar texto existente (placeholder)
                tf.clear() 

                points = slide_data.get('points', [])
                if points:
                    # Usar el primer párrafo (que clear() deja vacío pero existente)
                    p = tf.paragraphs[0]
                    p.text = points[0]
                    p.level = 0
                    p.space_after = Pt(10)
                    p.font.color.rgb = RGBColor(*theme['text'])
                    p.font.size = Pt(18)

                    # Añadir el resto
                    for point in points[1:]:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                        p.space_after = Pt(10)
                        p.font.color.rgb = RGBColor(*theme['text'])
                        p.font.size = Pt(18)

            # Guardar
            filename = f"presentation_{int(datetime.now().timestamp())}.pptx"
            local_path = os.path.join(self.output_folder, filename)
            prs.save(local_path)
            
            return f"{Config.PUBLIC_URL}/static/generated_docs/{filename}"
            
        except Exception as e:
            logger.error("Error PPTX: %s", e)
            return None

    def _generate_runway_video(self, prompt, input_image_path=None):
        try:
            logger.info("Llamando a Runway Gen-3 Alpha. Prompt: %s...", prompt[:50])

            # Preparar Imagen (Base64 Data URI)
            prompt_image_uri = None
            if input_image_path and os.path.exists(input_image_path):
                 import base64
                 import mimetypes
                 mime_type, _ = mimetypes.guess_type(input_image_path)
                 if not mime_type:
                     mime_type = "image/png"
                 
                 with open(input_image_path, "rb") as image_file:
                    encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
                    prompt_image_uri = f"data:{mime_type};base64,{encoded_string}"

            # 1. Iniciar Generación
            # Gen-3 Alpha Image-to-Video
            task_args = {
                "model": 'gen3a_turbo',
                "prompt_text": prompt,
            }
            if prompt_image_uri:
                task_args["prompt_image"] = prompt_image_uri
            
            task = self.runway_client.image_to_video.create(**task_args)
            
            task_id = task.id
            logger.info("Tarea Runway iniciada: %s", task_id)
            
            # 2. Polling con exponential backoff (max ~5 min acumulado)
            video_url = None
            wait = 3  # segundos iniciales
            elapsed = 0
            max_wait = 300  # timeout total
            _runway_start = time.time()

            while elapsed < max_wait:
                time.sleep(wait)
                elapsed += wait
                task_status = self.runway_client.tasks.retrieve(task_id)
                status = task_status.status
                logger.debug("Status Runway [%s]: %s (t+%ds)", task_id, status, elapsed)

                if status == "SUCCEEDED":
                    video_url = task_status.output[0]
                    _track(self.user_phone, "gen3a_turbo", "runway_video_generation",
                           latency_ms=int((time.time() - _runway_start) * 1000),
                           extra={"duration_s": 5})
                    break
                elif status == "FAILED":
                    _track(self.user_phone, "gen3a_turbo", "runway_video_generation",
                           latency_ms=int((time.time() - _runway_start) * 1000),
                           success=False, error=task_status.failure_reason)
                    logger.error("Runway task fallida [%s]: %s", task_id, task_status.failure_reason)
                    return None

                wait = min(wait * 2, 30)  # duplicar hasta un máximo de 30s

            if not video_url:
                return None

            # 3. Descargar Video
            filename = f"runway_{int(datetime.now().timestamp())}.mp4"
            local_path = os.path.join(self.output_folder, filename)
            
            r = requests.get(video_url)
            with open(local_path, 'wb') as f:
                f.write(r.content)
                
            logger.info("Video Runway descargado: %s", filename)
            return f"{Config.PUBLIC_URL}/static/generated_docs/{filename}"

        except Exception as e:
            logger.error("Error Runway API: %s", e)
            # Fallback a simulación si falla
            return self._generate_simulated_video(prompt)

    def _generate_simulated_video(self, prompt):
        logger.warning("Runway no disponible — RUNWAYML_API_SECRET no configurado o fallo de API.")
        return None

    def _apply_theme(self, slide, theme):
        # Fondo
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme['bg'])

    def _style_text(self, shape, rgb, bold=False, size=None):
        if not shape.text_frame:
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*rgb)
                if bold:
                    run.font.bold = True
                if size:
                    run.font.size = size

    def _refine_prompt_for_dalle(self, text, business_name):
        sys = """
        Eres un experto en Prompts para DALL-E 3.
        Objetivo: Crear una imagen publicitaria de altísimo impacto visual y estético.
        
        RESTRICCIONES DE TEXTO (STRICT):
        - NO incluyas frases largas, eslóganes ni oraciones.
        - SOLO se permite el Nombre de la Marca o Siglas.
        - El texto debe ser mínimo, elegante e integrado en el diseño.
        - Si duda, mejor SIN texto que con texto ilegible.
        
        INSTRUCCIÓN DE MARCA:
        - Incluye visiblemente la marca: '{business_name}'.
        """
        user_msg = f"Idea del usuario: {text}. Marca: {business_name}. Diseña algo moderno, minimalista y profesional."
        
        resp = self.openai.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "system", "content": sys}, {"role": "user", "content": user_msg}]
        )
        return resp.choices[0].message.content

    def _get_presentation_structure(self, topic):
        sys = """
        Eres un Consultor Senior de Negocio experto en storytelling.
        Tu objetivo es crear el GUIÓN DETALLADO para una presentación de alto impacto (5-6 diapositivas).
        
        INPUT: Tema del usuario (ej: "Consultoría Gemini").
        OUTPUT: JSON Estricto con el contenido real.
        
        REGLAS DE CONTENIDO:
        - NO uses frases vacías como "Punto 1". Escribe CONTENIDO REAL y VALIOSO.
        - Cada diapositiva debe tener entre 3 y 5 "points".
        - Los "points" deben ser frases completas y explicativas.
        
        REGLAS DE ESTILO (Theme):
        - Clasifica el tema en: 'TECH', 'ECO', 'CORPORATE', 'LUXURY', 'CREATIVE'.
        
        FORMATO JSON:
        {
            "theme": "TECH",
            "title": "TÍTULO PRINCIPAL ATRACTIVO",
            "subtitle": "Subtítulo que venda la idea",
            "slides": [
                {
                    "title": "1. El Desafío Actual",
                    "points": [
                        "Dato impactante: El 40% del tiempo se pierde en tareas repetitivas.",
                        "La competencia ya está adoptando IA generativa.",
                        "Coste de oportunidad de no innovar."
                    ]
                },
                ...
            ]
        }
        """
        try:
            resp = self.openai.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "system", "content": sys}, {"role": "user", "content": topic}],
                response_format={"type": "json_object"},
                max_tokens=1500
            )
            return json.loads(resp.choices[0].message.content)
        except Exception as e:
            logger.error("Error GPT JSON estructura presentación: %s", e)
            return {
                "theme": "CORPORATE",
                "title": "Error",
                "subtitle": "Intenta de nuevo",
                "slides": []
            }
