import os
import requests
import json
from datetime import datetime
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from runwayml import RunwayML

from core.config import Config

class MarketingAgent:
    THEMES = {
        "TECH": {"bg": (15, 23, 42), "title": (56, 189, 248), "text": (226, 232, 240)}, # Slate-900 / Sky-400 / Slate-200
        "ECO": {"bg": (20, 83, 45), "title": (74, 222, 128), "text": (220, 252, 231)},   # Green-900 / Green-400 / Green-100
        "CORPORATE": {"bg": (255, 255, 255), "title": (30, 58, 138), "text": (71, 85, 105)}, # White / Blue-900 / Slate-600
        "LUXURY": {"bg": (0, 0, 0), "title": (212, 175, 55), "text": (250, 250, 250)},   # Black / Gold / White
        "CREATIVE": {"bg": (255, 241, 242), "title": (190, 18, 60), "text": (88, 28, 135)} # Pink-50 / Rose-700 / Purple-900
    }

    def __init__(self):
        from core.clients import get_openai_client, get_runway_client
        self.openai = get_openai_client()
        self.runway_client = get_runway_client()
        # Configurar rutas de salida
        self.base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        self.output_folder = os.path.join(self.base_dir, 'static', 'generated_docs')
        os.makedirs(self.output_folder, exist_ok=True)

    def generate_marketing_content(self, prompt_text, content_type="image", business_name="Mi Empresa", logo_path=None):
        """
        Genera contenido creativo.
        content_type: 'image' (DALL-E 3) o 'slide' (PowerPoint).
        """
        print(f"🎨 MarketingAgent: Creando {content_type} sobre: '{prompt_text}' para {business_name}")
        
        if content_type == "image":
            return self._generate_dalle_image(prompt_text, business_name, logo_path)
        elif content_type == "slide":
            return self._generate_pptx_slide(prompt_text, business_name)
        elif content_type == "video":
            # 1. Visual Intelligence: Analizar Imagen para crear Prompt
            target_image = logo_path if logo_path else "dummy_product.jpg"
            video_prompt = self._analyze_product_context(target_image)
            
            # 2. Generar Video con Runway (Image + Text)
            # Pasamos la imagen original y el prompt mejorado
            video_url = self._generate_runway_video(video_prompt, input_image_path=target_image)
            return video_url
        else:
            return None

    def _generate_dalle_image(self, user_prompt, business_name, logo_path=None):
        try:
            # 1. Optimizar prompt para DALL-E usando GPT-4o
            optimized_prompt = self._refine_prompt_for_dalle(user_prompt, business_name)
            print(f"   -> Prompt DALL-E: {optimized_prompt[:50]}...")
            
            # 2. Generar Imagen
            response = self.openai.images.generate(
                model="dall-e-3",
                prompt=optimized_prompt,
                size="1024x1024",
                quality="standard",
                n=1,
            )
            image_url = response.data[0].url
            
            # 3. Descargar y guardar localmente
            filename = f"marketing_{int(datetime.now().timestamp())}.png"
            local_path = os.path.join(self.output_folder, filename)
            
            img_data = requests.get(image_url).content
            with open(local_path, 'wb') as handler:
                handler.write(img_data)
                
            # 4. Superponer Logo (Si existe)
            if logo_path:
                try:
                    from PIL import Image
                    # Convertir ruta web relativa a absoluta
                    # logo_path suele ser "/static/uploads/logos/..."
                    # self.base_dir está en "modules/proactive/...", subir 3 niveles para llegar a root
                    # Pero self.output_folder ya usa self.base_dir correctamente calculado en __init__
                    # self.base_dir en __init__ es: os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
                    # Esto apunta a TICKETIA_PRO root.
                    # logo_path empieza con /, quitarlo.
                    abs_logo_path = os.path.join(self.base_dir, logo_path.lstrip('/'))
                    
                    if os.path.exists(abs_logo_path):
                        base_img = Image.open(local_path).convert("RGBA")
                        logo_img = Image.open(abs_logo_path).convert("RGBA")
                        
                        # Resize logo (ej: 20% del ancho)
                        base_w, base_h = base_img.size
                        target_w = int(base_w * 0.20)
                        aspect_ratio = logo_img.width / logo_img.height
                        target_h = int(target_w / aspect_ratio)
                        logo_img = logo_img.resize((target_w, target_h), Image.Resampling.LANCZOS)
                        
                        # Posición: Esquina inferior derecha con padding
                        padding = 30
                        pos_x = base_w - target_w - padding
                        pos_y = base_h - target_h - padding
                        
                        # Pegar (usando el mismo logo como máscara para transparencia)
                        base_img.paste(logo_img, (pos_x, pos_y), logo_img)
                        
                        # Guardar (sobreescribir)
                        base_img.save(local_path, format="PNG")
                        print("✅ Logo superpuesto correctamente.")
                    else:
                        print(f"⚠️ Logo no encontrado en fs: {abs_logo_path}")
                except Exception as e_pil:
                    print(f"⚠️ Error procesando logo PIL: {e_pil}")
                
            return f"{Config.PUBLIC_URL}/static/generated_docs/{filename}"
            
        except Exception as e:
            print(f"❌ Error DALL-E: {e}")
            return None

    def _analyze_product_context(self, image_path):
        """
        Usa GPT-4o Vision para actuar como Director Creativo Universal.
        Ya no usamos if/else. La IA decide la estrategia visual según el objeto.
        """
        # Simulamos que le pasamos la imagen.
        # En producción real: encoded_image = encode_image(image_path)
        
        system_prompt = """
        You are an expert Prompt Engineer for Runway Gen-3. 
        The User is frustrated because the video outputs look like static images ("slideshow").
        
        YOUR GOAL: FORCE THE AI TO ANIMATE/TRANSFORM THE IMAGE.
        
        <AGGRESSIVE_STRATEGY>
        If the input is a static object (e.g. flat shirt), you must TRICK the video AI.
        Do NOT describe the image as it is. Describe what it MUST BECOME.
        State the "After" scene as a present fact.
        </AGGRESSIVE_STRATEGY>

        RULES FOR SPECIFIC OBJECTS:
        1. CLOTHING (Shirt/Dress): 
           - PROMPT: "Professional video of a fashion model WEARING this exact shirt walking on a beach. The fabric moves naturally with the body. Wind blowing."
           - KEYWORD: "WEARING" (Force the concept of a human).
        
        2. FOOD/DRINK:
           - PROMPT: "Slow motion commercial of this cake being SLICED by a silver knife. Crumbs falling. Steam rising."
           - KEYWORD: "SLICED" / "POURED" / "EATEN".
        
        3. TOOLS/MACHINES:
           - PROMPT: "Close up action shot of this drill SPINNING and boring into a wooden wall. Dust flying debris. Heavy vibration."
           - KEYWORD: "WORKING" / "SPINNING".

        4. STATIC OBJECTS (Perfume/Decor):
           - PROMPT: " The object floating in zero gravity, slowly rotating, water splashing around it, magical lighting."
        
        OUTPUT FORMAT (Direct & Punchy):
        "[SUBJECT IN ACTION] in [CONTEXT]. [MOVEMENT DETAILS]. Cinematic, photorealistic, 4k."

        Examples:
        - "A fashion model wearing this floral shirt walking confidently on a sunny boardwalk, wind blowing the fabric."
        - "A silver knife slicing through this cheesecake, revealing the creamy texture, crumbs falling."
        - "This electric drill spinning rapidly and drilling a hole in a wall, sawdust flying everywhere."

        Be BOLD. Demand MOVEMENT.
        Responde SOLO con el prompt final en inglés.
        """

        try:
            # 1. Preparar imagen (Base64)
            import base64
            import mimetypes
            
            if not os.path.exists(image_path):
                print("❌ Imagen no encontrada para análisis.")
                return "Cinematic shot of the product, 4k"

            mime_type, _ = mimetypes.guess_type(image_path)
            if not mime_type: mime_type = "image/png"
            
            with open(image_path, "rb") as image_file:
                encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
                base64_image = f"data:{mime_type};base64,{encoded_string}"
            
            # 2. Llamada REAL a GPT-4o
            print(f"👁️ Visual Intelligence: Analizando {os.path.basename(image_path)} con GPT-4o...")
            
            response = self.openai.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system", 
                        "content": system_prompt
                    },
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "This is the product image. Generate the perfect video prompt."},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": base64_image,
                                    "detail": "high"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=150
            )

            prompt_generated = response.choices[0].message.content.strip()
            print(f"🧠 Prompt Generado: {prompt_generated}")
            return prompt_generated
                
        except Exception as e:
            print(f"❌ Error Visual Intelligence: {e}")
            # Fallback seguro
            return "Cinematic rotating studio shot of the product on a podium, dramatic rim lighting, minimalist background, 4k"
                
        except Exception as e:
            print(f"❌ Error en análisis visual: {e}")
            return "Cinematic product shot, 4k, professional lighting"
    def _generate_pptx_slide(self, topic, business_name):
        try:
            # 1. Estructurar contenido COMPLETO + TEMA
            presentation_data = self._get_presentation_structure(topic)
            theme_key = presentation_data.get('theme', 'CORPORATE').upper()
            theme = self.THEMES.get(theme_key, self.THEMES['CORPORATE'])
            
            # 2. Crear PPT
            prs = Presentation()
            
            # --- Slide 1: Título Principal ---
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            self._apply_theme(slide, theme) # Aplicar estilos
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = presentation_data.get('title', 'Propuesta')
            # Incluir nombre empresa en subtítulo
            subtitle.text = f"{presentation_data.get('subtitle', '')}\nPropuesto por: {business_name}"
            
            self._style_text(title, theme['title'], bold=True, size=Pt(44))
            self._style_text(subtitle, theme['text'], size=Pt(24))

            # --- Slides de Contenido ---
            for slide_data in presentation_data.get('slides', []):
                bullet_slide_layout = prs.slide_layouts[1] # Título + Bullets
                slide = prs.slides.add_slide(bullet_slide_layout)
                self._apply_theme(slide, theme)
                
                # Título de la Slide
                shapes = slide.shapes
                title_shape = shapes.title
                title_shape.text = slide_data.get('title', 'Detalle')
                self._style_text(title_shape, theme['title'], bold=True, size=Pt(36))
                
                # Cuerpo (Bullets)
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.word_wrap = True
                # Limpiar texto existente (placeholder)
                tf.clear() 

                points = slide_data.get('points', [])
                if points:
                    # Usar el primer párrafo (que clear() deja vacío pero existente)
                    p = tf.paragraphs[0]
                    p.text = points[0]
                    p.level = 0
                    p.space_after = Pt(10)
                    p.font.color.rgb = RGBColor(*theme['text'])
                    p.font.size = Pt(18)

                    # Añadir el resto
                    for point in points[1:]:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                        p.space_after = Pt(10)
                        p.font.color.rgb = RGBColor(*theme['text'])
                        p.font.size = Pt(18)

            # Guardar
            filename = f"presentation_{int(datetime.now().timestamp())}.pptx"
            local_path = os.path.join(self.output_folder, filename)
            prs.save(local_path)
            
            return f"{Config.PUBLIC_URL}/static/generated_docs/{filename}"
            
        except Exception as e:
            print(f"❌ Error PPTX: {e}")
            return None

    def _generate_runway_video(self, prompt, input_image_path=None):
        try:
            print(f"🎬 MarketingAgent: Llamando a Runway Gen-3 Alpha...")
            print(f"   -> Prompt: {prompt[:50]}...")
            print(f"   -> Input Image: {input_image_path}")

            # Preparar Imagen (Base64 Data URI)
            prompt_image_uri = None
            if input_image_path and os.path.exists(input_image_path):
                 import base64
                 import mimetypes
                 mime_type, _ = mimetypes.guess_type(input_image_path)
                 if not mime_type: mime_type = "image/png"
                 
                 with open(input_image_path, "rb") as image_file:
                    encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
                    prompt_image_uri = f"data:{mime_type};base64,{encoded_string}"

            # 1. Iniciar Generación
            # Gen-3 Alpha Image-to-Video
            task_args = {
                "model": 'gen3a_turbo',
                "prompt_text": prompt,
            }
            if prompt_image_uri:
                task_args["prompt_image"] = prompt_image_uri
            
            task = self.runway_client.image_to_video.create(**task_args)
            
            task_id = task.id
            print(f"   -> Tarea Runway iniciada: {task_id}")
            
            # 2. Polling de estado (Esperar hasta 3 min)
            import time
            status = "PENDING"
            video_url = None
            
            # Polling loop con timeout de veiligheid
            for _ in range(60): # Max 5 mins (60 * 5s)
                time.sleep(5) 
                task_status = self.runway_client.tasks.retrieve(task_id)
                status = task_status.status
                print(f"   -> Status Runway: {status}")
                
                if status == "SUCCEEDED":
                    video_url = task_status.output[0] 
                    break
                elif status == "FAILED":
                    print(f"❌ Runway Task Failed: {task_status.failure_reason}")
                    return None
            
            if not video_url: return None

            # 3. Descargar Video
            filename = f"runway_{int(datetime.now().timestamp())}.mp4"
            local_path = os.path.join(self.output_folder, filename)
            
            r = requests.get(video_url)
            with open(local_path, 'wb') as f:
                f.write(r.content)
                
            print("✅ Video Runway descargado correctamente.")
            return f"{Config.PUBLIC_URL}/static/generated_docs/{filename}"
                
        except Exception as e:
            print(f"❌ Error Runway API: {e}")
            # Fallback a simulación si falla
            return self._generate_simulated_video(prompt)

    def _generate_simulated_video(self, prompt):
        # This is a dummy function for demonstration purposes.
        # In a real scenario, this would generate a placeholder video or return a static URL.
        print(f"   -> Generando video simulado para: {prompt[:50]}...")
        # Simulate a video URL
        return f"{Config.PUBLIC_URL}/static/generated_docs/simulated_video.mp4"

    def _apply_theme(self, slide, theme):
        # Fondo
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme['bg'])

    def _style_text(self, shape, rgb, bold=False, size=None):
        if not shape.text_frame: return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*rgb)
                if bold: run.font.bold = True
                if size: run.font.size = size

    def _refine_prompt_for_dalle(self, text, business_name):
        sys = """
        Eres un experto en Prompts para DALL-E 3.
        Objetivo: Crear una imagen publicitaria de altísimo impacto visual y estético.
        
        RESTRICCIONES DE TEXTO (STRICT):
        - NO incluyas frases largas, eslóganes ni oraciones.
        - SOLO se permite el Nombre de la Marca o Siglas.
        - El texto debe ser mínimo, elegante e integrado en el diseño.
        - Si duda, mejor SIN texto que con texto ilegible.
        
        INSTRUCCIÓN DE MARCA:
        - Incluye visiblemente la marca: '{business_name}'.
        """
        user_msg = f"Idea del usuario: {text}. Marca: {business_name}. Diseña algo moderno, minimalista y profesional."
        
        resp = self.openai.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "system", "content": sys}, {"role": "user", "content": user_msg}]
        )
        return resp.choices[0].message.content

    def _get_presentation_structure(self, topic):
        sys = """
        Eres un Consultor Senior de Negocio experto en storytelling.
        Tu objetivo es crear el GUIÓN DETALLADO para una presentación de alto impacto (5-6 diapositivas).
        
        INPUT: Tema del usuario (ej: "Consultoría Gemini").
        OUTPUT: JSON Estricto con el contenido real.
        
        REGLAS DE CONTENIDO:
        - NO uses frases vacías como "Punto 1". Escribe CONTENIDO REAL y VALIOSO.
        - Cada diapositiva debe tener entre 3 y 5 "points".
        - Los "points" deben ser frases completas y explicativas.
        
        REGLAS DE ESTILO (Theme):
        - Clasifica el tema en: 'TECH', 'ECO', 'CORPORATE', 'LUXURY', 'CREATIVE'.
        
        FORMATO JSON:
        {
            "theme": "TECH",
            "title": "TÍTULO PRINCIPAL ATRACTIVO",
            "subtitle": "Subtítulo que venda la idea",
            "slides": [
                {
                    "title": "1. El Desafío Actual",
                    "points": [
                        "Dato impactante: El 40% del tiempo se pierde en tareas repetitivas.",
                        "La competencia ya está adoptando IA generativa.",
                        "Coste de oportunidad de no innovar."
                    ]
                },
                ...
            ]
        }
        """
        try:
            resp = self.openai.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "system", "content": sys}, {"role": "user", "content": topic}],
                response_format={"type": "json_object"},
                max_tokens=1500
            )
            return json.loads(resp.choices[0].message.content)
        except Exception as e:
            print(f"Error GPT JSON: {e}")
            return {
                "theme": "CORPORATE",
                "title": "Error",
                "subtitle": "Intenta de nuevo",
                "slides": []
            }
